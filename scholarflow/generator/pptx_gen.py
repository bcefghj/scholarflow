"""Generate PowerPoint presentation with embedded figures and smart layouts."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu

from scholarflow.models import PaperContent, PaperFigure, SlideData

THEMES = {
    "academic_blue": {
        "bg": RGBColor(0xFF, 0xFF, 0xFF),
        "title_color": RGBColor(0x1A, 0x3C, 0x6E),
        "text_color": RGBColor(0x33, 0x33, 0x33),
        "accent": RGBColor(0x2E, 0x75, 0xB6),
        "footer_bg": RGBColor(0x1A, 0x3C, 0x6E),
        "light_bg": RGBColor(0xF0, 0xF4, 0xF8),
    },
    "minimal_white": {
        "bg": RGBColor(0xFA, 0xFA, 0xFA),
        "title_color": RGBColor(0x22, 0x22, 0x22),
        "text_color": RGBColor(0x44, 0x44, 0x44),
        "accent": RGBColor(0xE8, 0x4D, 0x39),
        "footer_bg": RGBColor(0x22, 0x22, 0x22),
        "light_bg": RGBColor(0xF5, 0xF5, 0xF5),
    },
    "dark_modern": {
        "bg": RGBColor(0x1E, 0x1E, 0x2E),
        "title_color": RGBColor(0xCD, 0xD6, 0xF4),
        "text_color": RGBColor(0xBA, 0xC2, 0xDE),
        "accent": RGBColor(0x89, 0xB4, 0xFA),
        "footer_bg": RGBColor(0x11, 0x11, 0x1B),
        "light_bg": RGBColor(0x2A, 0x2A, 0x3C),
    },
}


def generate_pptx(
    slides: list[SlideData],
    content: PaperContent,
    output_dir: Path,
    theme: str = "academic_blue",
) -> Path:
    """Generate a .pptx presentation file with embedded figures."""
    output_dir.mkdir(parents=True, exist_ok=True)
    colors = THEMES.get(theme, THEMES["academic_blue"])

    # Build figure lookup: fig_id -> PaperFigure
    fig_lookup: dict[str, PaperFigure] = {}
    for fig in content.figures:
        fig_lookup[fig.figure_id] = fig

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for i, slide_data in enumerate(slides):
        if i == 0:
            _add_title_slide(prs, slide_data, content, colors)
        else:
            fig = _find_figure_for_slide(slide_data, fig_lookup)
            if fig and fig.path.exists():
                _add_figure_slide(prs, slide_data, fig, colors, i)
            else:
                _add_content_slide(prs, slide_data, colors, i)

    out_path = output_dir / "slides.pptx"
    prs.save(str(out_path))
    return out_path


def _find_figure_for_slide(slide_data: SlideData, fig_lookup: dict[str, PaperFigure]):
    """Check if a slide references a figure via figure_path field."""
    if not slide_data.figure_path:
        return None
    fig_id = slide_data.figure_path.lower().strip()
    if fig_id in fig_lookup:
        return fig_lookup[fig_id]
    for fid, fig in fig_lookup.items():
        if fid in fig_id or fig_id in fid:
            return fig
    return None


def _add_title_slide(prs, slide_data, content, colors):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = colors["footer_bg"]

    title_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11), Inches(2))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = content.meta.title or slide_data.title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.alignment = PP_ALIGN.CENTER

    info_box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(11), Inches(1.5))
    tf2 = info_box.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = ", ".join(content.meta.authors) if content.meta.authors else ""
    p2.font.size = Pt(18)
    p2.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
    p2.alignment = PP_ALIGN.CENTER

    if content.meta.year or content.meta.venue:
        p3 = tf2.add_paragraph()
        parts = [x for x in [content.meta.venue, content.meta.year] if x]
        p3.text = " | ".join(parts)
        p3.font.size = Pt(16)
        p3.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
        p3.alignment = PP_ALIGN.CENTER


def _add_figure_slide(prs, slide_data, fig: PaperFigure, colors, slide_num):
    """Add a slide with left text + right figure layout."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = colors["bg"]

    # Accent bar
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.15), Inches(7.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = colors["accent"]
    bar.line.fill.background()

    # Title (full width)
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = slide_data.title
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = colors["title_color"]

    # Left: bullet text (60% width)
    text_left = Inches(0.6)
    text_top = Inches(1.4)
    text_width = Inches(6.5)
    text_height = Inches(5.5)

    if slide_data.bullets:
        body_box = slide.shapes.add_textbox(text_left, text_top, text_width, text_height)
        tf = body_box.text_frame
        tf.word_wrap = True
        for j, bullet in enumerate(slide_data.bullets):
            para = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            para.text = f"  {bullet}"
            para.font.size = Pt(16)
            para.font.color.rgb = colors["text_color"]
            para.space_after = Pt(10)

    # Right: figure (40% width)
    img_left = Inches(7.5)
    img_top = Inches(1.4)
    img_max_w = Inches(5.3)
    img_max_h = Inches(4.5)

    try:
        pic = slide.shapes.add_picture(str(fig.path), img_left, img_top)
        # Scale to fit
        w_ratio = img_max_w / pic.width if pic.width > img_max_w else 1.0
        h_ratio = img_max_h / pic.height if pic.height > img_max_h else 1.0
        ratio = min(w_ratio, h_ratio)
        if ratio < 1.0:
            pic.width = int(pic.width * ratio)
            pic.height = int(pic.height * ratio)
        # Center horizontally in the right zone
        pic.left = int(img_left + (img_max_w - pic.width) / 2)
    except Exception:
        pass

    # Caption below image
    if fig.caption:
        cap_box = slide.shapes.add_textbox(img_left, Inches(6.2), img_max_w, Inches(0.8))
        tf = cap_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = fig.caption[:120]
        p.font.size = Pt(10)
        p.font.italic = True
        p.font.color.rgb = RGBColor(0x88, 0x88, 0x88)
        p.alignment = PP_ALIGN.CENTER

    # Page number
    _add_page_number(slide, colors, slide_num)


def _add_content_slide(prs, slide_data, colors, slide_num):
    """Add a text-only content slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = colors["bg"]

    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.15), Inches(7.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = colors["accent"]
    bar.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12), Inches(0.9))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = slide_data.title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = colors["title_color"]

    if slide_data.bullets:
        body_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.5), Inches(5))
        tf = body_box.text_frame
        tf.word_wrap = True
        for j, bullet in enumerate(slide_data.bullets):
            para = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            para.text = f"  {bullet}"
            para.font.size = Pt(18)
            para.font.color.rgb = colors["text_color"]
            para.space_after = Pt(12)

    _add_page_number(slide, colors, slide_num)


def _add_page_number(slide, colors, slide_num):
    page_box = slide.shapes.add_textbox(Inches(12), Inches(7), Inches(1), Inches(0.4))
    tf = page_box.text_frame
    p = tf.paragraphs[0]
    p.text = str(slide_num + 1)
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
    p.alignment = PP_ALIGN.RIGHT
